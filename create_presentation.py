from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Traffic Speed Prediction Using LSTM"
subtitle.text = "A project presentation on predicting traffic speed using LSTM neural networks"

# Slide 2: Project Overview
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]

title.text = "Project Overview"
content.text = ("- Objective: Predict traffic speed at various nodes using LSTM\n"
                "- Background: Importance of traffic prediction for smart cities\n"
                "- Data: Node ID, Speed, Timestamp")

# Slide 3: Data Preprocessing
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]

title.text = "Data Preprocessing"
content.text = ("- Load and inspect data\n"
                "- Handle missing values and convert data types\n"
                "- Save processed data by Node ID\n")

# Slide 4: Exploratory Data Analysis (EDA)
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]

title.text = "Exploratory Data Analysis (EDA)"
content.text = ("- Statistical summary of data\n"
                "- Time series visualization\n"
                "- Speed distribution analysis by Node ID\n")

# Slide 5: LSTM Model Building and Training
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]

title.text = "LSTM Model Building and Training"
content.text = ("- Overview of LSTM model architecture\n"
                "- Split data into training and testing sets\n"
                "- Train the model and save the trained model\n")

# Slide 6: Model Evaluation and Prediction
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.placeholders[1]

title.text = "Model Evaluation and Prediction"
content.text = ("- Compare predicted data with actual data\n"
                "- Visualize prediction results\n"
                "- Save results for future analysis\n")

# Slide 7: Conclusion and Future Work
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.placeholders[1]

title.text = "Conclusion and Future Work"
content.text = ("- Summary of project and key findings\n"
                "- Interpretation of results\n"
                "- Future improvements and additional tasks\n")

# Save the presentation
prs.save("Traffic_Speed_Prediction_Presentation.pptx")
